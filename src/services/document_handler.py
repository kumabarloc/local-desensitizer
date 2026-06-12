"""文档格式解析与生成"""
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Optional
import re


class DocumentHandler(ABC):
    """文档处理器基类"""

    @abstractmethod
    def read(self, path: Path) -> str:
        """读取文档内容（返回纯文本）"""
        pass

    @abstractmethod
    def write(self, path: Path, content: str) -> None:
        """写入文档（纯文本或对应格式）"""
        pass

    @abstractmethod
    def get_extension(self) -> str:
        """返回文件扩展名"""
        pass


class TextHandler(DocumentHandler):
    """纯文本处理器 .txt"""

    def read(self, path: Path) -> str:
        with open(path, 'r', encoding='utf-8') as f:
            return f.read()

    def write(self, path: Path, content: str) -> None:
        with open(path, 'w', encoding='utf-8') as f:
            f.write(content)

    def get_extension(self) -> str:
        return ".txt"


class MarkdownHandler(DocumentHandler):
    """Markdown处理器 .md"""

    def read(self, path: Path) -> str:
        with open(path, 'r', encoding='utf-8') as f:
            return f.read()

    def write(self, path: Path, content: str) -> None:
        with open(path, 'w', encoding='utf-8') as f:
            f.write(content)

    def get_extension(self) -> str:
        return ".md"


class CsvHandler(DocumentHandler):
    """CSV处理器 .csv"""

    def read(self, path: Path) -> str:
        # 读取CSV返回纯文本内容
        import csv
        lines = []
        with open(path, 'r', encoding='utf-8-sig') as f:
            reader = csv.reader(f)
            for row in reader:
                lines.append(','.join(row))
        return '\n'.join(lines)

    def write(self, path: Path, content: str) -> None:
        import csv
        lines = content.split('\n')
        with open(path, 'w', encoding='utf-8-sig', newline='') as f:
            writer = csv.writer(f)
            for line in lines:
                writer.writerow(line.split(','))

    def get_extension(self) -> str:
        return ".csv"


class DocxHandler(DocumentHandler):
    """Word文档处理器 .docx

    读写都走 Markdown 中间格式，以保留表格/标题/列表结构。
    - read(): 用 python-docx 按 body 顺序遍历，输出 MD 语法
            （不走 mammoth，因为 mammoth 转 MD 会丢表格）
    - write(): MD → docx (markdown_to_docx)
    """

    def read(self, path: Path) -> str:
        from docx import Document
        from docx.oxml.ns import qn

        doc = Document(path)
        # 建立 XML 元素到 python-docx 对象的映射（O(n) 查找）
        para_map = {p._p: p for p in doc.paragraphs}
        table_map = {t._tbl: t for t in doc.tables}

        parts: list[str] = []
        body = doc.element.body
        for child in body.iterchildren():
            if child.tag == qn('w:p'):
                if child in para_map:
                    p = para_map[child]
                    md_line = self._paragraph_to_md(p)
                    parts.append(md_line)
            elif child.tag == qn('w:tbl'):
                if child in table_map:
                    t = table_map[child]
                    md_table = self._table_to_md(t)
                    parts.append(md_table)

        # 页眉页脚
        extra: list[str] = []
        for section in doc.sections:
            header = section.header
            if header:
                for para in header.paragraphs:
                    if para.text.strip():
                        extra.append(f"[Header] {para.text.strip()}")
            footer = section.footer
            if footer:
                for para in footer.paragraphs:
                    if para.text.strip():
                        extra.append(f"[Footer] {para.text.strip()}")
        if extra:
            parts.append('\n'.join(extra))

        return '\n\n'.join(parts)

    @staticmethod
    def _paragraph_to_md(p) -> str:
        """段落 → MD 语法"""
        text = p.text
        if not text.strip():
            return ''  # 空行

        style = p.style.name if p.style else 'Normal'

        # 标题
        if 'Heading' in style or style == 'Title':
            level = 1
            for k in ('Heading 1', 'Title'):
                if style == k:
                    return f'# {text}'
            for lvl in range(1, 7):
                if style == f'Heading {lvl}':
                    return f'{'#' * lvl} {text}'
            return f'## {text}'  # fallback

        # 列表
        if 'List' in style and 'Bullet' in style.replace('List ', ''):
            return f'- {text}'
        if 'List' in style and ('Number' in style or 'Paragraph' in style):
            return f'1. {text}'  # 简化为有序起点，docx 不强制数字连续

        # 普通段落
        return text

    @staticmethod
    def _table_to_md(t) -> str:
        """表格 → MD 语法"""
        rows = []
        for row in t.rows:
            cells = []
            for cell in row.cells:
                # 单元格可能含多行，用空格拼接避免破坏 MD 表格
                cell_text = cell.text.strip().replace('\n', ' ')
                cells.append(cell_text)
            rows.append('| ' + ' | '.join(cells) + ' |')

        # 插入表头分隔行
        if rows:
            n_cols = rows[0].count('|') - 1
            if n_cols > 0:
                separator = '| ' + ' | '.join(['---'] * n_cols) + ' |'
                rows.insert(1, separator)

        return '\n'.join(rows)

    def write(self, path: Path, content: str) -> None:
        from src.services.md_converter import markdown_to_docx
        markdown_to_docx(content, path)

    def get_extension(self) -> str:
        return ".docx"


class XlsxHandler(DocumentHandler):
    """Excel处理器 .xlsx"""

    def read(self, path: Path) -> str:
        from openpyxl import load_workbook
        wb = load_workbook(path, data_only=True)
        lines = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            lines.append(f"=== {sheet_name} ===")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else '' for c in row]
                if any(c.strip() for c in cells):
                    lines.append('\t'.join(cells))
        return '\n'.join(lines)

    def write(self, path: Path, content: str) -> None:
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        for i, line in enumerate(content.split('\n'), start=1):
            ws.cell(row=i, column=1, value=line)
        wb.save(path)

    def get_extension(self) -> str:
        return ".xlsx"


class PptxHandler(DocumentHandler):
    """PowerPoint处理器 .pptx"""

    def read(self, path: Path) -> str:
        from pptx import Presentation
        prs = Presentation(path)
        lines = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            lines.append(f"--- Slide {slide_num} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip())
        return '\n'.join(lines)

    def write(self, path: Path, content: str) -> None:
        from pptx import Presentation
        from pptx.util import Pt
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        for i, line in enumerate(content.split('\n')):
            txBox = slide.shapes.add_textbox(Pt(50), Pt(50 + i * 20), Pt(600), Pt(20))
            tf = txBox.text_frame
            tf.text = line
        prs.save(path)

    def get_extension(self) -> str:
        return ".pptx"


def get_handler(ext: str) -> Optional[DocumentHandler]:
    """根据扩展名获取对应处理器"""
    handlers = {
        '.txt': TextHandler(),
        '.md': MarkdownHandler(),
        '.csv': CsvHandler(),
        '.docx': DocxHandler(),
        '.xlsx': XlsxHandler(),
        '.pptx': PptxHandler(),
    }
    return handlers.get(ext.lower())


def detect_encoding(filepath: Path) -> str:
    """检测文件编码（主要用于CSV）"""
    import chardet
    with open(filepath, 'rb') as f:
        raw = f.read(10000)
    result = chardet.detect(raw)
    return result['encoding'] or 'utf-8'


def read_document(filepath: Path) -> tuple[str, DocumentHandler]:
    """读取任意支持的文档，返回 (纯文本内容, 处理器)"""
    ext = filepath.suffix
    handler = get_handler(ext)
    if not handler:
        raise ValueError(f"不支持的文档格式: {ext}")
    content = handler.read(filepath)
    return content, handler


def write_document(filepath: Path, content: str, original_ext: str) -> None:
    """写入文档，保持原始格式"""
    ext = filepath.suffix or original_ext
    handler = get_handler(ext)
    if not handler:
        raise ValueError(f"不支持的文档格式: {ext}")
    handler.write(filepath, content)